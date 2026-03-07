import fitz  # PyMuPDF
import docx
from pptx import Presentation


def parse_and_chunk(file_stream, filename, chunk_size=500):
    ext = filename.lower().split(".")[-1]
    if ext == "pdf":
        return chunk_pdf(file_stream, chunk_size)
    elif ext == "docx":
        return chunk_docx(file_stream, chunk_size)
    elif ext == "pptx":
        return chunk_pptx(file_stream, chunk_size)
    elif ext == "txt":
        return chunk_txt(file_stream, chunk_size)
    else:
        raise ValueError("Unsupported file type")


def chunk_pdf(file_stream, chunk_size):
    file_stream.seek(0)
    doc = fitz.open(stream=file_stream.read(), filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return [text[i : i + chunk_size] for i in range(0, len(text), chunk_size)]


def chunk_docx(file_stream, chunk_size):
    file_stream.seek(0)
    doc = docx.Document(file_stream)
    text = "\n".join([p.text for p in doc.paragraphs])
    return [text[i : i + chunk_size] for i in range(0, len(text), chunk_size)]


def chunk_pptx(file_stream, chunk_size):
    file_stream.seek(0)
    prs = Presentation(file_stream)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return [text[i : i + chunk_size] for i in range(0, len(text), chunk_size)]


def chunk_txt(file_stream, chunk_size):
    file_stream.seek(0)
    text = file_stream.read().decode("utf-8")
    return [text[i : i + chunk_size] for i in range(0, len(text), chunk_size)]
