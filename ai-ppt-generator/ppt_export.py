from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os
import uuid
import uuid
from visual_generator import get_visual_for_slide

def hex_to_rgb(hex_color):
    """Converts hex color string to RGBColor object."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def apply_background(slide, theme):
    """Applies the theme background color to the slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    try:
        fill.fore_color.rgb = hex_to_rgb(theme.get('background_color', '#FFFFFF'))
    except:
        fill.fore_color.rgb = RGBColor(255, 255, 255)

def add_dynamic_design(slide, theme):
    """Adds generative design elements based on theme."""
    # Add a dynamic accent bar or shape
    accent_color = theme.get('accent_color', '#000000')
    
    # Example: A stylish top bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.2))
    shape.fill.solid()
    try:
        shape.fill.fore_color.rgb = hex_to_rgb(accent_color)
    except:
        pass
    shape.line.fill.background()

def create_slide(prs, slide_data, theme):
    # Use blank layout to paint from scratch
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 1. Apply Theme
    apply_background(slide, theme)
    add_dynamic_design(slide, theme)
    
    # Colors
    title_color = hex_to_rgb(theme.get('primary_color', '#000000'))
    text_color = hex_to_rgb(theme.get('text_color', '#333333'))
    
    # 2. Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1.5))
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = slide_data.get('title', 'Untitled')
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = theme.get('font_family', 'Arial')
    
    # 3. Visuals & Layout
    layout = slide_data.get('layout', 'content_text_left')
    visual_hint = slide_data.get('visual_description')
    image_path = None
    
    if visual_hint:
        try:
            image_path = get_visual_for_slide(visual_hint)
        except:
            pass
            
    # Layout Logic
    if layout == 'title':
        # Centered Title
        title_shape.top = Inches(2.5)
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(60)
        
        # Subtitle
        content = slide_data.get('content', [])
        if content:
            subtitle = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(2))
            sp = subtitle.text_frame.paragraphs[0]
            sp.text = content[0] if isinstance(content, list) else content
            sp.alignment = PP_ALIGN.CENTER
            sp.font.size = Pt(28)
            sp.font.color.rgb = text_color
            
    elif image_path and os.path.exists(image_path) and layout != 'full_text':
        # Split Layout
        if layout == 'content_text_right':
            # Image Left, Text Right
            pic = slide.shapes.add_picture(image_path, Inches(0.5), Inches(2), height=Inches(5))
            content_left = Inches(7)
        else:
            # Text Left, Image Right (Standard)
            content_left = Inches(0.5)
            pic = slide.shapes.add_picture(image_path, Inches(7.5), Inches(2), height=Inches(5))
            
        # Content Box
        body = slide.shapes.add_textbox(content_left, Inches(2), Inches(6), Inches(5))
        tf = body.text_frame
        tf.word_wrap = True
        
        for point in slide_data.get('content', []):
            bp = tf.add_paragraph()
            bp.text = point
            bp.font.size = Pt(20)
            bp.font.color.rgb = text_color
            bp.space_after = Pt(14)
            bp.level = 0
            
    else:
        # Full Width Content
        body = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(5))
        tf = body.text_frame
        tf.word_wrap = True
        
        for point in slide_data.get('content', []):
            bp = tf.add_paragraph()
            bp.text = point
            bp.font.size = Pt(24)
            bp.font.color.rgb = text_color
            bp.space_after = Pt(18)

    # Speaker Notes
    if 'speaker_notes' in slide_data:
        slide.notes_slide.notes_text_frame.text = slide_data['speaker_notes']


def export_to_pptx(data, out_file=None):
    """
    Exports the presentation. 
    data can be:
      - A list of slides (legacy support)
      - A dict {'theme': {...}, 'slides': [...]} (new dynamic design)
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    if isinstance(data, dict) and 'slides' in data:
        slides = data['slides']
        theme = data.get('theme', {})
    else:
        slides = data
        theme = {}
        
    for slide_data in slides:
        create_slide(prs, slide_data, theme)
            
    # Save file
    pptx_dir = 'generated_pptx'
    os.makedirs(pptx_dir, exist_ok=True)
    if not out_file:
        out_file = f"CurioSlides_{uuid.uuid4().hex[:8]}.pptx"
    out_path = os.path.join(pptx_dir, out_file)
    prs.save(out_path)
    return out_file
