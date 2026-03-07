from langchain_community.vectorstores import FAISS
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
import fitz
import docx
from pptx import Presentation
import os
from dotenv import load_dotenv
import uuid

load_dotenv()
GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')

print("Using Gemini embeddings: models/text-embedding-004")
embedding = GoogleGenerativeAIEmbeddings(
    model="models/text-embedding-004", 
    google_api_key=GEMINI_API_KEY
)

def extract_text(file_stream, filename):
    ext = filename.lower().split('.')[-1]
    file_stream.seek(0)
    
    if ext == 'pdf':
        try:
            # fitz.open with stream expects bytes if filetype is specified
            doc = fitz.open(stream=file_stream.read(), filetype='pdf')
            return "\n".join([page.get_text() for page in doc])
        except Exception as e:
            raise ValueError(f"Error parsing PDF: {e}")
            
    elif ext == 'docx':
        try:
            doc = docx.Document(file_stream)
            return "\n".join([p.text for p in doc.paragraphs])
        except Exception as e:
            raise ValueError(f"Error parsing DOCX: {e}")
            
    elif ext == 'pptx':
        try:
            prs = Presentation(file_stream)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except Exception as e:
            raise ValueError(f"Error parsing PPTX: {e}")
            
    elif ext == 'txt':
        try:
            return file_stream.read().decode('utf-8')
        except Exception as e:
            raise ValueError(f"Error parsing TXT: {e}")
            
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def store_document(file_stream, filename):
    doc_id = str(uuid.uuid4())
    raw_text = extract_text(file_stream, filename)
    
    if not raw_text.strip():
        raise ValueError("No text extracted from document")

    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    docs = splitter.create_documents([raw_text])

    # Attach doc_id and filename as metadata to each chunk
    for doc in docs:
        doc.metadata = getattr(doc, 'metadata', {})
        doc.metadata['doc_id'] = doc_id
        doc.metadata['filename'] = filename

    db_path = "vector_db"
    
    try:
        if os.path.exists(os.path.join(db_path, "index.faiss")):
            vectordb = FAISS.load_local(db_path, embedding, allow_dangerous_deserialization=True)
            vectordb.add_documents(docs)
        else:
            vectordb = FAISS.from_documents(docs, embedding)
        
        vectordb.save_local(db_path)
    except Exception as e:
        raise ValueError(f"Error updating vector store: {e}")
        
    return doc_id
